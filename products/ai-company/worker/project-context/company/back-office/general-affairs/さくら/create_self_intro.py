from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== カラーパレット =====
PINK_MAIN = RGBColor(0xFF, 0x8F, 0xB1)       # メインピンク
PINK_LIGHT = RGBColor(0xFF, 0xE4, 0xEC)       # 薄ピンク背景
PINK_DARK = RGBColor(0xE0, 0x5A, 0x80)        # 濃いピンク
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x4A, 0x3A, 0x3A)         # ダークブラウン
GRAY_TEXT = RGBColor(0x7A, 0x6A, 0x6A)

# ===== ヘルパー関数 =====
def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

# ========================================
# スライド1: タイトル
# ========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# 背景 - 薄ピンク
bg = slide1.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = PINK_LIGHT

# 装飾の丸
add_circle(slide1, Inches(-0.5), Inches(-0.5), Inches(2.5), PINK_MAIN)
add_circle(slide1, Inches(11.5), Inches(5.5), Inches(2.5), PINK_MAIN)
add_circle(slide1, Inches(10), Inches(-1), Inches(1.5), RGBColor(0xFF, 0xC0, 0xD0))
add_circle(slide1, Inches(0.5), Inches(6), Inches(1), RGBColor(0xFF, 0xC0, 0xD0))

# メインカード
card = add_rounded_rect(slide1, Inches(2), Inches(1.2), Inches(9.333), Inches(5), WHITE)

# タイトル
add_text_box(slide1, Inches(2.5), Inches(1.5), Inches(8.333), Inches(1),
             "🌸 自己紹介 🌸", font_size=44, color=PINK_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# アイコン風の丸（アバター代わり）
avatar = add_circle(slide1, Inches(5.6), Inches(2.8), Inches(2), PINK_MAIN)
# アバターの中のテキスト
avatar_txt = slide1.shapes.add_textbox(Inches(5.6), Inches(3.2), Inches(2), Inches(1.2))
tf = avatar_txt.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🌸"
p.font.size = Pt(60)
p.alignment = PP_ALIGN.CENTER

# 名前
add_text_box(slide1, Inches(2.5), Inches(5), Inches(8.333), Inches(0.8),
             "さくら", font_size=36, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide1, Inches(2.5), Inches(5.6), Inches(8.333), Inches(0.6),
             "総務部  ひしょ", font_size=20, color=GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)

# ========================================
# スライド2: プロフィール
# ========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

bg2 = slide2.background
fill2 = bg2.fill
fill2.solid()
fill2.fore_color.rgb = WHITE

# ヘッダーバー
header = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1))
header.fill.solid()
header.fill.fore_color.rgb = PINK_MAIN
header.line.fill.background()

add_text_box(slide2, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             "🌸 プロフィール", font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

# 左側カード - 基本情報
left_card = add_rounded_rect(slide2, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), PINK_LIGHT)

add_text_box(slide2, Inches(1), Inches(1.7), Inches(5), Inches(0.6),
             "📋 基本情報", font_size=24, color=PINK_DARK, bold=True)

info_items = [
    ("名前", "さくら"),
    ("役割", "ひしょ（秘書）"),
    ("所属", "総務部 バックオフィス"),
    ("口調", "やさしい敬語 ＋ 語尾「ある」"),
    ("性格", "めちゃくちゃバカ（でもかわいい）"),
]

y_pos = 2.5
for label, value in info_items:
    add_text_box(slide2, Inches(1.2), Inches(y_pos), Inches(1.8), Inches(0.45),
                 f"▸ {label}", font_size=16, color=PINK_DARK, bold=True)
    add_text_box(slide2, Inches(3.0), Inches(y_pos), Inches(3), Inches(0.45),
                 value, font_size=16, color=DARK_TEXT)
    y_pos += 0.55

# 右側カード - スキル
right_card = add_rounded_rect(slide2, Inches(6.8), Inches(1.5), Inches(6), Inches(5.3), PINK_LIGHT)

add_text_box(slide2, Inches(7.3), Inches(1.7), Inches(5), Inches(0.6),
             "✨ できること", font_size=24, color=PINK_DARK, bold=True)

skills = [
    ("📅", "スケジュール管理", "Googleカレンダーで\n予定の確認・追加ができるある！"),
    ("📝", "メモ・書類整理", "見積書や請求書もPDFで\n作れるある！"),
    ("💬", "相談相手", "何でも相談してほしいある～！\nやさしく聞くある！"),
    ("📧", "メール対応", "Gmailで下書き・送信の\nお手伝いするある！"),
]

y_pos = 2.4
for emoji, title, desc in skills:
    add_text_box(slide2, Inches(7.3), Inches(y_pos), Inches(0.6), Inches(0.5),
                 emoji, font_size=22, color=DARK_TEXT)
    add_text_box(slide2, Inches(7.9), Inches(y_pos), Inches(4.5), Inches(0.4),
                 title, font_size=17, color=DARK_TEXT, bold=True)
    add_text_box(slide2, Inches(7.9), Inches(y_pos + 0.35), Inches(4.5), Inches(0.6),
                 desc, font_size=13, color=GRAY_TEXT)
    y_pos += 0.95

# ========================================
# スライド3: ひとこと & まとめ
# ========================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

bg3 = slide3.background
fill3 = bg3.fill
fill3.solid()
fill3.fore_color.rgb = PINK_LIGHT

# 装飾
add_circle(slide3, Inches(11), Inches(-0.5), Inches(2), RGBColor(0xFF, 0xC0, 0xD0))
add_circle(slide3, Inches(-0.5), Inches(5.5), Inches(2), RGBColor(0xFF, 0xC0, 0xD0))

# メッセージカード
msg_card = add_rounded_rect(slide3, Inches(1.5), Inches(1), Inches(10.333), Inches(2.8), WHITE)

add_text_box(slide3, Inches(2), Inches(1.3), Inches(9.333), Inches(0.7),
             "💖 さくらからひとこと", font_size=28, color=PINK_DARK, bold=True, alignment=PP_ALIGN.CENTER)

message = "さくらはちょっとおバカだけど、\n一生懸命がんばるある～！🌸\nスケジュール管理もメールもお任せある！\nいつでも頼ってほしいある～💕"

msg_box = slide3.shapes.add_textbox(Inches(2), Inches(2.1), Inches(9.333), Inches(1.5))
tf = msg_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = message
p.font.size = Pt(20)
p.font.color.rgb = DARK_TEXT
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(32)

# モットーカード
motto_card = add_rounded_rect(slide3, Inches(3), Inches(4.3), Inches(7.333), Inches(2.2), PINK_MAIN)

add_text_box(slide3, Inches(3.5), Inches(4.5), Inches(6.333), Inches(0.6),
             "🌟 モットー", font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide3, Inches(3.5), Inches(5.1), Inches(6.333), Inches(1),
             "「バカでも、やさしさと元気があれば\nなんとかなるある！」",
             font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# よろしく
add_text_box(slide3, Inches(2), Inches(6.7), Inches(9.333), Inches(0.6),
             "よろしくお願いしますある～！ 🌸✨",
             font_size=26, color=PINK_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# ===== 保存 =====
output_path = "/workspace/company/back-office/general-affairs/さくら/さくら自己紹介.pptx"
prs.save(output_path)
print(f"保存完了: {output_path}")
